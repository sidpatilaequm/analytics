"""exporters.py — turn one report's computed results into a PDF, a PPTX, or
an XLSX. All three read the exact same normalized structure (see
`app.py:_render_boxes`), so a box that's wrong in one format is wrong (and
fixable) in exactly one place, not three.
"""

import os
import io
import re
from pptx.enum.shapes import MSO_SHAPE


# --------------------------------------------------------------------------
# BROWSER PDF
# --------------------------------------------------------------------------

# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------
def build_browser_pdf(url):
    """
    Generate a PDF from the actual rendered published React report.

    The dashboard is rendered using normal screen CSS so the grid,
    box sizes and alignment remain identical to the published page.

    The rendered dashboard is captured at high resolution and then
    placed into an A4 PDF at 300 DPI for sharp text and charts.
    """

    from playwright.sync_api import sync_playwright
    from PIL import Image
    import io

    with sync_playwright() as p:
        browser = p.chromium.launch(
            headless=True,
            args=[
                "--no-sandbox",
                "--disable-setuid-sandbox",
            ],
        )

        # Keep the same desktop layout as the published dashboard.
        # Higher device scale factor = higher resolution screenshot.
        page = browser.new_page(
            viewport={
                "width": 1440,
                "height": 900,
            },
            device_scale_factor=3,
        )

        response = page.goto(
            url,
            wait_until="networkidle",
        )

        print("PDF BROWSER URL:", url)
        print(
            "PDF PAGE STATUS:",
            response.status if response else None
        )
        print("PDF PAGE URL:", page.url)

        page.wait_for_selector(
            ".sheet",
            timeout=30000,
        )

        # Wait for React/charts/fonts.
        page.wait_for_timeout(1500)

        page.evaluate(
            """
            async () => {
                if (document.fonts) {
                    await document.fonts.ready;
                }
            }
            """
        )

        # Hide download buttons only.
        page.add_style_tag(
            content="""
                .export-actions {
                    display: none !important;
                }
            """
        )

        # --------------------------------------------------------------
        # Capture the EXACT rendered dashboard.
        # --------------------------------------------------------------

        sheet = page.locator(".sheet")

        dimensions = sheet.evaluate(
            """
            el => {
                const r = el.getBoundingClientRect();

                return {
                    x: r.x,
                    y: r.y,
                    width: r.width,
                    height: r.height
                };
            }
            """
        )

        print("PDF SHEET CSS DIMENSIONS:", dimensions)

        screenshot = sheet.screenshot(
            type="png",
            animations="disabled",
        )

        # Close browser only after screenshot is captured.
        browser.close()

    # --------------------------------------------------------------
    # Convert screenshot to high-resolution A4 pages.
    #
    # A4 @ 300 DPI:
    #   8.27in × 300 = 2480 px
    #   11.69in × 300 = 3508 px
    # --------------------------------------------------------------

    image = Image.open(
        io.BytesIO(screenshot)
    ).convert("RGB")

    print("PDF SCREENSHOT PIXELS:", image.size)

    A4_WIDTH = 2480
    A4_HEIGHT = 3508

    # Preserve exact aspect ratio.
    scale = A4_WIDTH / image.width

    scaled_width = A4_WIDTH
    scaled_height = round(image.height * scale)

    image = image.resize(
        (scaled_width, scaled_height),
        Image.Resampling.LANCZOS,
    )

    print(
        "PDF SCALED IMAGE:",
        image.size
    )

    # --------------------------------------------------------------
    # Split dashboard into A4 pages.
    # --------------------------------------------------------------

    pages = []

    for top in range(0, scaled_height, A4_HEIGHT):

        bottom = min(
            top + A4_HEIGHT,
            scaled_height,
        )

        page_image = Image.new(
            "RGB",
            (A4_WIDTH, A4_HEIGHT),
            "white",
        )

        crop = image.crop(
            (
                0,
                top,
                scaled_width,
                bottom,
            )
        )

        page_image.paste(
            crop,
            (0, 0),
        )

        pages.append(page_image)

    # --------------------------------------------------------------
    # Generate PDF at 300 DPI.
    # --------------------------------------------------------------

    output = io.BytesIO()

    pages[0].save(
        output,
        format="PDF",
        resolution=300.0,
        save_all=True,
        append_images=pages[1:],
    )

    return output.getvalue()


def _esc(s):
    return (str(s) if s is not None else "").replace("&", "&amp;") \
        .replace("<", "&lt;").replace(">", "&gt;")


def _cell(v):
    return "—" if v is None else str(v)


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------
def build_pptx(name, sections):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(40)
    p.font.bold = True

    for sec in sections:
        slide = prs.slides.add_slide(blank)
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
        tp = title.text_frame.paragraphs[0]
        tp.text = sec["name"]
        tp.font.size = Pt(28)
        tp.font.bold = True

        y = 1.3
        for box in sec["boxes"]:
            if y > 6.6:
                slide = prs.slides.add_slide(blank)
                y = 0.5

            if box.get("error"):
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12), Inches(0.6))
                tb.text_frame.text = f"{box['title']}: {box['error']}"
                y += 0.7
                continue

            if box["kind"] == "value":

                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5),
                    Inches(y),
                    Inches(5.8),
                    Inches(1.5),
                )

                # Card background
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(255, 255, 255)

                # Card border
                card.line.color.rgb = RGBColor(213, 221, 218)

                tf = card.text_frame
                tf.clear()
                tf.word_wrap = True

                # Title
                p1 = tf.paragraphs[0]
                p1.text = box["title"] or ""
                p1.font.size = Pt(12)
                p1.font.color.rgb = RGBColor(0x5A, 0x66, 0x63)

                # Main value
                p2 = tf.add_paragraph()
                p2.text = box.get("text", "—")
                p2.font.size = Pt(30)
                p2.font.bold = True

                # Optional note
                if box.get("note"):
                    p3 = tf.add_paragraph()
                    p3.text = box["note"]
                    p3.font.size = Pt(8)
                    p3.font.color.rgb = RGBColor(120, 120, 120)

                y += 1.7

            elif box["kind"] in ("chart", "table") and box.get("rows"):
                rows = box["rows"][:13]  # keep one slide readable
                r, c = len(rows), max(len(x) for x in rows)
                label = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(6), Inches(0.4))
                label.text_frame.text = box["title"] or ""
                label.text_frame.paragraphs[0].font.bold = True
                y += 0.45
                height = min(0.32 * r, 4.8)
                shape = slide.shapes.add_table(r, c, Inches(0.5), Inches(y),
                                                Inches(12.3), Inches(height))
                tbl = shape.table
                for ri in range(r):
                    for ci in range(c):
                        v = rows[ri][ci] if ci < len(rows[ri]) else None
                        cell = tbl.cell(ri, ci)
                        cell.text = "—" if v is None else str(v)
                        cell.text_frame.paragraphs[0].font.size = Pt(11)
                y += height + 0.35

            elif box.get("text"):
                nb = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(1))
                nb.text_frame.word_wrap = True
                nb.text_frame.text = f"{box['title']}: {box['text']}" if box["title"] else box["text"]
                y += 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------
# XLSX
# --------------------------------------------------------------------------
def build_xlsx(name, sections):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    wb.remove(wb.active)
    used = set()

    for sec in sections:
        base = re.sub(r"[\[\]\*\?/\\:]", "", sec["name"] or "Section")[:28] or "Section"
        sheet_name, i = base, 2
        while sheet_name in used:
            sheet_name = f"{base[:26]}-{i}"
            i += 1
        used.add(sheet_name)
        ws = wb.create_sheet(sheet_name)

        row = 1
        for box in sec["boxes"]:
            ws.cell(row=row, column=1, value=box["title"] or "Untitled").font = Font(bold=True)
            row += 1
            if box.get("error"):
                ws.cell(row=row, column=1, value=f"Error: {box['error']}")
                row += 2
                continue
            if box["kind"] == "value":
                ws.cell(row=row, column=1, value=box.get("text", "—"))
                if box.get("note"):
                    ws.cell(row=row, column=2, value=box["note"])
                row += 2
            elif box["kind"] in ("chart", "table") and box.get("rows"):
                for ri, r in enumerate(box["rows"]):
                    for ci, v in enumerate(r):
                        cell = ws.cell(row=row + ri, column=1 + ci, value=v)
                        if ri == 0:
                            cell.font = Font(bold=True)
                            cell.fill = PatternFill("solid", fgColor="E7ECEA")
                row += len(box["rows"]) + 2
            elif box.get("text"):
                ws.cell(row=row, column=1, value=box["text"])
                row += 2
        for col in range(1, 9):
            ws.column_dimensions[chr(64 + col)].width = 22

    if not wb.sheetnames:
        wb.create_sheet("Report")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()

# --------------------------------------------------------------------------
# CSV / TABLE
# --------------------------------------------------------------------------
def build_csv(name, sections):
    import csv

    buf = io.StringIO()
    writer = csv.writer(buf)

    writer.writerow([name])

    for sec in sections:
        writer.writerow([])
        writer.writerow([sec["name"]])

        if sec.get("desc"):
            writer.writerow([sec["desc"]])

        for box in sec["boxes"]:
            writer.writerow([])
            writer.writerow([box["title"] or "Untitled"])

            if box.get("error"):
                writer.writerow(["Error", box["error"]])

            elif box["kind"] == "value":
                writer.writerow(["Value", box.get("text", "—")])

                if box.get("note"):
                    writer.writerow(["Note", box["note"]])

            elif box["kind"] in ("chart", "table") and box.get("rows"):
                for row in box["rows"]:
                    writer.writerow([
                        "—" if value is None else value
                        for value in row
                    ])

            elif box.get("text"):
                writer.writerow([box["text"]])

    return buf.getvalue().encode("utf-8")
# --------------------------------------------------------------------------
# TABLE PDF
# --------------------------------------------------------------------------
def build_table_pdf(name, sections):
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
        PageBreak,
    )

    buf = io.BytesIO()

    doc = SimpleDocTemplate(
        buf,
        pagesize=landscape(A4),
        topMargin=36,
        bottomMargin=36,
        leftMargin=36,
        rightMargin=36,
    )

    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "TableReportTitle",
        parent=styles["Title"],
        fontSize=20,
        spaceAfter=20,
    )

    section_style = ParagraphStyle(
        "TableSection",
        parent=styles["Heading2"],
        fontSize=14,
        spaceBefore=14,
        spaceAfter=8,
    )

    header_style = ParagraphStyle(
        "TableHeader",
        parent=styles["BodyText"],
        fontSize=9,
        fontName="Helvetica-Bold",
        textColor=colors.white,
    )

    cell_style = ParagraphStyle(
        "TableCell",
        parent=styles["BodyText"],
        fontSize=9,
        leading=12,
    )

    story = [
        Paragraph(_esc(name), title_style),
    ]

    for sec in sections:

        # Section heading
        story.append(
            Paragraph(
                _esc(sec.get("name", "Untitled Section")),
                section_style,
            )
        )

        if sec.get("desc"):
            story.append(
                Paragraph(
                    _esc(sec["desc"]),
                    cell_style,
                )
            )
            story.append(Spacer(1, 8))

        # Main table header
        table_data = [
            [
                Paragraph("Metric", header_style),
                Paragraph("Value", header_style),
                Paragraph("Note / Details", header_style),
            ]
        ]

        for box in sec.get("boxes", []):

            # Error boxes
            if box.get("error"):
                table_data.append([
                    Paragraph(_esc(box.get("title") or "Error"), cell_style),
                    Paragraph("Error", cell_style),
                    Paragraph(_esc(box["error"]), cell_style),
                ])

            # Value / KPI boxes
            elif box.get("kind") == "value":
                table_data.append([
                    Paragraph(
                        _esc(box.get("title") or "Untitled"),
                        cell_style,
                    ),
                    Paragraph(
                        _esc(box.get("text", "—")),
                        cell_style,
                    ),
                    Paragraph(
                        _esc(box.get("note", "")),
                        cell_style,
                    ),
                ])

            # Chart or table boxes
            elif (
                box.get("kind") in ("chart", "table")
                and box.get("rows")
            ):
                rows = box["rows"]

                # Add the box title
                table_data.append([
                    Paragraph(
                        f"<b>{_esc(box.get('title') or 'Table')}</b>",
                        cell_style,
                    ),
                    "",
                    "",
                ])

                # Add every row from the chart/table
                for row in rows:
                    row_values = [
                        "—" if value is None else str(value)
                        for value in row
                    ]

                    if len(row_values) == 1:
                        table_data.append([
                            Paragraph(_esc(row_values[0]), cell_style),
                            "",
                            "",
                        ])

                    elif len(row_values) == 2:
                        table_data.append([
                            Paragraph(_esc(row_values[0]), cell_style),
                            Paragraph(_esc(row_values[1]), cell_style),
                            "",
                        ])

                    else:
                        table_data.append([
                            Paragraph(_esc(row_values[0]), cell_style),
                            Paragraph(_esc(row_values[1]), cell_style),
                            Paragraph(
                                _esc(" | ".join(row_values[2:])),
                                cell_style,
                            ),
                        ])

            # Notes
            elif box.get("text"):
                table_data.append([
                    Paragraph(
                        _esc(box.get("title") or "Details"),
                        cell_style,
                    ),
                    Paragraph(
                        _esc(box.get("text", "")),
                        cell_style,
                    ),
                    "",
                ])

        # Create section table
        if len(table_data) > 1:
            table = Table(
                table_data,
                colWidths=[
                    3.2 * inch,
                    2.0 * inch,
                    4.0 * inch,
                ],
                repeatRows=1,
                hAlign="LEFT",
            )

            table.setStyle(TableStyle([
                # Header
                (
                    "BACKGROUND",
                    (0, 0),
                    (-1, 0),
                    colors.HexColor("#0D6E62"),
                ),
                (
                    "TEXTCOLOR",
                    (0, 0),
                    (-1, 0),
                    colors.white,
                ),

                # Grid
                (
                    "GRID",
                    (0, 0),
                    (-1, -1),
                    0.5,
                    colors.HexColor("#D5DDDA"),
                ),

                # Alternating rows
                (
                    "BACKGROUND",
                    (0, 1),
                    (-1, -1),
                    colors.white,
                ),

                # Alignment
                (
                    "VALIGN",
                    (0, 0),
                    (-1, -1),
                    "MIDDLE",
                ),

                # Padding
                (
                    "TOPPADDING",
                    (0, 0),
                    (-1, -1),
                    8,
                ),
                (
                    "BOTTOMPADDING",
                    (0, 0),
                    (-1, -1),
                    8,
                ),
                (
                    "LEFTPADDING",
                    (0, 0),
                    (-1, -1),
                    10,
                ),
                (
                    "RIGHTPADDING",
                    (0, 0),
                    (-1, -1),
                    10,
                ),

                # Value column centered
                (
                    "ALIGN",
                    (1, 1),
                    (1, -1),
                    "CENTER",
                ),
            ]))

            story.append(table)
            story.append(Spacer(1, 16))

    doc.build(story)

    return buf.getvalue()